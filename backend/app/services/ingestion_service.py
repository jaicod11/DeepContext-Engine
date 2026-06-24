"""
services/ingestion_service.py
------------------------------
Full document ingestion pipeline:
  load → split → embed → upsert to Pinecone

Page-aware chunking for PDF, PPTX, XLSX:
  - PDF:  chunks tagged with page_number (1-indexed)
  - PPTX: chunks tagged with slide_number (1-indexed)
  - XLSX: chunks tagged with sheet_name
  - Other formats: page_number = None
"""

from __future__ import annotations

import hashlib
from dataclasses import dataclass
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter

from app.core.config import Settings, get_settings
from app.core.logging import get_logger
from app.vectorstore.embeddings import CachedEmbedder, get_embedder
from app.vectorstore.pinecone_client import PineconeClient, VectorRecord, get_pinecone_client

logger = get_logger(__name__)


@dataclass
class IngestionResult:
    document_id:      str
    filename:         str
    chunks_total:     int
    vectors_upserted: int
    namespace:        str


class IngestionService:
    def __init__(
        self,
        pinecone: PineconeClient | None = None,
        embedder: CachedEmbedder | None = None,
        settings: Settings | None = None,
    ) -> None:
        self._pc       = pinecone or get_pinecone_client()
        self._embedder = embedder or get_embedder()
        self._settings = settings or get_settings()
        self._splitter = RecursiveCharacterTextSplitter(
            chunk_size=self._settings.chunk_size,
            chunk_overlap=self._settings.chunk_overlap,
            separators=self._settings.splitter_separators,
        )

    # ─────────────────────────────────────────
    # Public API
    # ─────────────────────────────────────────

    async def ingest_file(
        self,
        file_path:   str | Path,
        namespace:   str | None = None,
        metadata:    dict | None = None,
        document_id: str | None = None,
    ) -> IngestionResult:
        path   = Path(file_path)
        ns     = namespace or self._settings.pinecone_namespace
        doc_id = document_id or self._file_hash(path)

        logger.info("ingestion_started", file=path.name, document_id=doc_id, namespace=ns)

        # Load page-aware content: list of (chunk_text, extra_metadata)
        page_chunks = await self._load_paged(path)

        # Delete stale vectors if re-ingesting
        await self._pc.delete_by_filter(
            filter={"document_id": {"$eq": doc_id}},
            namespace=ns,
        )

        # Embed all chunks
        texts   = [chunk for chunk, _ in page_chunks]
        vectors = await self._embedder.embed_documents(texts)

        # Build VectorRecord list — each chunk carries its page/slide/sheet metadata
        base_meta = {
            "document_id":  doc_id,
            "source":       path.name,
            "total_chunks": len(page_chunks),
            **(metadata or {}),
        }
        records = [
            VectorRecord(
                vector_id=f"{doc_id}-{i}",
                values=vec,
                metadata={
                    **base_meta,
                    **chunk_meta,          # page_number, slide_number, or sheet_name
                    "text":        chunk,
                    "chunk_index": i,
                },
            )
            for i, ((chunk, chunk_meta), vec) in enumerate(zip(page_chunks, vectors))
        ]

        upserted = await self._pc.upsert(records, namespace=ns)

        logger.info(
            "ingestion_complete",
            file=path.name,
            document_id=doc_id,
            chunks=len(page_chunks),
            upserted=upserted,
            namespace=ns,
        )
        return IngestionResult(
            document_id=doc_id,
            filename=path.name,
            chunks_total=len(page_chunks),
            vectors_upserted=upserted,
            namespace=ns,
        )

    async def ingest_text(
        self,
        text:        str,
        source:      str = "inline",
        namespace:   str | None = None,
        metadata:    dict | None = None,
        document_id: str | None = None,
    ) -> IngestionResult:
        ns     = namespace or self._settings.pinecone_namespace
        doc_id = document_id or hashlib.sha256(text.encode()).hexdigest()[:16]

        chunks  = self._splitter.split_text(text)
        vectors = await self._embedder.embed_documents(chunks)

        await self._pc.delete_by_filter(
            filter={"document_id": {"$eq": doc_id}},
            namespace=ns,
        )

        records = [
            VectorRecord(
                vector_id=f"{doc_id}-{i}",
                values=vec,
                metadata={
                    "document_id":  doc_id,
                    "source":       source,
                    "text":         chunk,
                    "chunk_index":  i,
                    "total_chunks": len(chunks),
                    "page_number":  None,
                    **(metadata or {}),
                },
            )
            for i, (chunk, vec) in enumerate(zip(chunks, vectors))
        ]

        upserted = await self._pc.upsert(records, namespace=ns)
        return IngestionResult(
            document_id=doc_id,
            filename=source,
            chunks_total=len(chunks),
            vectors_upserted=upserted,
            namespace=ns,
        )

    async def delete_document(
        self,
        document_id: str,
        namespace:   str | None = None,
    ) -> None:
        ns = namespace or self._settings.pinecone_namespace
        await self._pc.delete_by_filter(
            filter={"document_id": {"$eq": document_id}},
            namespace=ns,
        )
        logger.info("document_deleted", document_id=document_id, namespace=ns)

    # ─────────────────────────────────────────
    # Page-aware loaders
    # Returns: list of (chunk_text, extra_metadata_dict)
    # ─────────────────────────────────────────

    async def _load_paged(self, path: Path) -> list[tuple[str, dict]]:
        suffix = path.suffix.lower()

        if suffix == ".pdf":
            return await self._load_pdf_paged(path)
        elif suffix == ".pptx":
            return await self._load_pptx_paged(path)
        elif suffix in (".xlsx", ".xls"):
            return await self._load_xlsx_paged(path)
        elif suffix == ".docx":
            text = await self._load_docx(path)
        elif suffix in (".txt", ".md"):
            text = await self._load_text(path)
        elif suffix in (".html", ".htm"):
            text = await self._load_html(path)
        else:
            raise ValueError(
                f"Unsupported file type: {suffix}. "
                f"Supported: .pdf, .docx, .txt, .md, .html, .pptx, .xlsx, .xls"
            )

        # For non-page-aware formats, split normally with no page metadata
        chunks = self._splitter.split_text(text)
        return [(chunk, {"page_number": None}) for chunk in chunks]

    async def _load_pdf_paged(self, path: Path) -> list[tuple[str, dict]]:
        """
        Extract PDF text page by page.
        Each chunk is tagged with its 1-indexed page_number.
        Splitting happens within each page so page attribution is exact.
        """
        import asyncio
        from pypdf import PdfReader

        def _read() -> list[tuple[int, str]]:
            reader = PdfReader(str(path))
            pages = []
            for i, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append((i, text))
            return pages

        raw_pages = await asyncio.to_thread(_read)

        result: list[tuple[str, dict]] = []
        for page_num, page_text in raw_pages:
            chunks = self._splitter.split_text(page_text)
            for chunk in chunks:
                result.append((chunk, {"page_number": page_num}))

        return result

    async def _load_pptx_paged(self, path: Path) -> list[tuple[str, dict]]:
        """
        Extract PPTX text slide by slide.
        Each chunk is tagged with its 1-indexed slide_number.
        """
        import asyncio
        from pptx import Presentation

        def _read() -> list[tuple[int, str]]:
            prs = Presentation(str(path))
            slides = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                parts = [f"Slide {slide_num}"]
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes] {notes}")
                if len(parts) > 1:
                    slides.append((slide_num, "\n".join(parts)))
            return slides

        raw_slides = await asyncio.to_thread(_read)

        result: list[tuple[str, dict]] = []
        for slide_num, slide_text in raw_slides:
            chunks = self._splitter.split_text(slide_text)
            for chunk in chunks:
                result.append((chunk, {"slide_number": slide_num, "page_number": slide_num}))
        return result

    async def _load_xlsx_paged(self, path: Path) -> list[tuple[str, dict]]:
        """
        Extract XLSX text sheet by sheet.
        Each chunk is tagged with its sheet_name.
        """
        import asyncio
        import openpyxl

        def _read() -> list[tuple[str, str]]:
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            sheets = []
            for sheet_name in wb.sheetnames:
                ws  = wb[sheet_name]
                rows = [f"Sheet: {sheet_name}"]
                for row in ws.iter_rows(values_only=True):
                    if not any(cell is not None for cell in row):
                        continue
                    rows.append("\t".join(str(c) if c is not None else "" for c in row))
                if len(rows) > 1:
                    sheets.append((sheet_name, "\n".join(rows)))
            wb.close()
            return sheets

        raw_sheets = await asyncio.to_thread(_read)

        result: list[tuple[str, dict]] = []
        for sheet_name, sheet_text in raw_sheets:
            chunks = self._splitter.split_text(sheet_text)
            for chunk in chunks:
                result.append((chunk, {"sheet_name": sheet_name, "page_number": None}))
        return result

    # ─────────────────────────────────────────
    # Non-paged loaders (return full text)
    # ─────────────────────────────────────────

    @staticmethod
    async def _load_docx(path: Path) -> str:
        import asyncio
        from docx import Document

        def _read() -> str:
            doc = Document(str(path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

        return await asyncio.to_thread(_read)

    @staticmethod
    async def _load_text(path: Path) -> str:
        import aiofiles
        async with aiofiles.open(path, encoding="utf-8", errors="replace") as f:
            return await f.read()

    @staticmethod
    async def _load_html(path: Path) -> str:
        import asyncio
        import aiofiles
        from html.parser import HTMLParser

        class _Stripper(HTMLParser):
            def __init__(self) -> None:
                super().__init__()
                self._parts: list[str] = []
            def handle_data(self, data: str) -> None:
                stripped = data.strip()
                if stripped:
                    self._parts.append(stripped)
            def get_text(self) -> str:
                return " ".join(self._parts)

        async with aiofiles.open(path, encoding="utf-8", errors="replace") as f:
            html = await f.read()

        def _strip(html: str) -> str:
            p = _Stripper()
            p.feed(html)
            return p.get_text()

        return await asyncio.to_thread(_strip, html)

    # ─────────────────────────────────────────
    # Utilities
    # ─────────────────────────────────────────

    @staticmethod
    def _file_hash(path: Path) -> str:
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(65536), b""):
                h.update(chunk)
        return h.hexdigest()[:16]


# ─────────────────────────────────────────────
# Singleton
# ─────────────────────────────────────────────

_ingestion_instance: IngestionService | None = None


def get_ingestion_service() -> IngestionService:
    global _ingestion_instance
    if _ingestion_instance is None:
        _ingestion_instance = IngestionService()
    return _ingestion_instance